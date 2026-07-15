"""Build the native PowerPoint version of the defense deck.

    python presentation/build_pptx.py

Mirrors presentation/deck_src.html slide for slide (17 slides, 16:9): same
copy, same design system (indigo accent, stat tiles, for-the-experts asides),
result figures embedded as pictures, diagrams rebuilt as native editable
PowerPoint shapes. Output: presentation/final_presentation.pptx.
"""
from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
FIGS = ROOT / "results" / "figures"
OUT = Path(__file__).with_name("final_presentation.pptx")

# ---- design tokens (same as the HTML deck) ---------------------------------
ACCENT = RGBColor(0x4F, 0x46, 0xE5)
ACCENT_DARK = RGBColor(0x37, 0x30, 0xA3)
CYAN = RGBColor(0x08, 0x91, 0xB2)
ACCENT_SOFT = RGBColor(0xEE, 0xF0, 0xFD)
INK = RGBColor(0x1C, 0x21, 0x28)
MUTED = RGBColor(0x5B, 0x64, 0x70)
FAINT = RGBColor(0x8A, 0x93, 0xA2)
GOOD = RGBColor(0x16, 0xA3, 0x4A)
BAD = RGBColor(0xDC, 0x26, 0x26)
BORDER = RGBColor(0xE7, 0xE9, 0xF0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OFFW = RGBColor(0xFB, 0xFB, 0xFD)
LILAC = RGBColor(0xDF, 0xE3, 0xFF)
FONT = "Segoe UI"

SW, SH = Inches(13.333), Inches(7.5)
MX = Inches(0.65)          # side margin


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _noshadow(shape):
    shape.shadow.inherit = False
    return shape


def rect(slide, x, y, w, h, fill=WHITE, line=BORDER, radius=0.10):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    try:
        sp.adjustments[0] = radius
    except Exception:
        pass
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(1)
    return _noshadow(sp)


def text(slide, x, y, w, h, paras, *, size=18, color=INK, bold=False,
         align=PP_ALIGN.LEFT, spacing=1.15, anchor=MSO_ANCHOR.TOP,
         space_after=4):
    """paras: list of paragraphs; each is a str or a list of (txt, overrides) runs."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        p.space_after = Pt(space_after)
        runs = [(para, {})] if isinstance(para, str) else para
        for txt, ov in runs:
            r = p.add_run()
            r.text = txt
            r.font.name = FONT
            r.font.size = Pt(ov.get("size", size))
            r.font.bold = ov.get("bold", bold)
            r.font.italic = ov.get("italic", False)
            r.font.color.rgb = ov.get("color", color)
    return tb


def kicker(slide, txt, color=ACCENT, y=Inches(0.42)):
    text(slide, MX, y, SW - 2 * MX, Inches(0.32), [txt.upper()],
         size=13, color=color, bold=True)


def h2(slide, txt, y=Inches(0.72)):
    text(slide, MX, y, SW - 2 * MX, Inches(0.62), [txt], size=30, bold=True)


def footer(slide, n):
    text(slide, MX, Inches(7.08), Inches(6), Inches(0.3),
         ["CESI × KMUTNB — Research Initiation 2026"], size=10.5, color=FAINT)
    text(slide, SW - Inches(1.6), Inches(7.08), Inches(1.0), Inches(0.3),
         [str(n)], size=10.5, color=FAINT, align=PP_ALIGN.RIGHT)


def tile(slide, x, y, w, h, number, label, num_color=INK):
    rect(slide, x, y, w, h)
    text(slide, x + Inches(0.18), y + Inches(0.10), w - Inches(0.36), Inches(0.62),
         [number], size=30, bold=True, color=num_color)
    text(slide, x + Inches(0.18), y + Inches(0.68), w - Inches(0.36), h - Inches(0.74),
         [label], size=12, color=MUTED, spacing=1.1)


def card(slide, x, y, w, h, fill=OFFW, line=BORDER):
    return rect(slide, x, y, w, h, fill=fill, line=line)


def expert(slide, x, y, w, h, body):
    rect(slide, x, y, w, h, fill=ACCENT_SOFT, line=None, radius=0.08)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Pt(3.2), h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    _noshadow(bar)
    text(slide, x + Inches(0.18), y + Inches(0.06), w - Inches(0.3), Inches(0.24),
         ["FOR THE EXPERTS"], size=9.5, color=ACCENT, bold=True)
    text(slide, x + Inches(0.18), y + Inches(0.30), w - Inches(0.3), h - Inches(0.36),
         [body], size=12, color=RGBColor(0x3B, 0x3F, 0x6E), spacing=1.15)


def pic_fit(slide, name, x, y, w, h, framed=True):
    path = FIGS / f"{name}.png"
    iw, ih = Image.open(path).size
    scale = min(w / iw, h / ih)
    pw, ph = int(iw * scale), int(ih * scale)
    px, py = int(x + (w - pw) / 2), int(y + (h - ph) / 2)
    if framed:
        rect(slide, x, y, w, h, fill=WHITE, line=BORDER, radius=0.04)
    slide.shapes.add_picture(str(path), Emu(px), Emu(py), Emu(pw), Emu(ph))


def _arrowhead(line_shape, color):
    ln = line_shape.line._get_or_add_ln()
    tail = ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"})
    ln.append(tail)
    line_shape.line.color.rgb = color
    line_shape.line.width = Pt(2)


def arrow(slide, x1, y1, x2, y2, color=FAINT):
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    _arrowhead(ln, color)
    return _noshadow(ln)


def bullet(slide, x, y, w, items, size=15, gap=0.52, color=INK, mark=ACCENT):
    for i, it in enumerate(items):
        yy = y + Inches(gap * i)
        sq = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, yy + Inches(0.075),
                                    Inches(0.12), Inches(0.12))
        sq.fill.solid(); sq.fill.fore_color.rgb = mark; sq.line.fill.background()
        _noshadow(sq)
        runs = [(it, {})] if isinstance(it, str) else it
        text(slide, x + Inches(0.26), yy - Inches(0.03), w - Inches(0.26), Inches(gap),
             [runs], size=size, color=color, spacing=1.1)


def gradient_bg(slide):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    sp.line.fill.background()
    _noshadow(sp)
    try:
        sp.fill.gradient()
        stops = sp.fill.gradient_stops
        stops[0].color.rgb = ACCENT_DARK
        stops[0].position = 0.0
        stops[1].color.rgb = CYAN
        stops[1].position = 1.0
        try:
            sp.fill.gradient_angle = 35.0
        except Exception:
            pass
    except Exception:
        sp.fill.solid()
        sp.fill.fore_color.rgb = ACCENT
    return sp


# ============================ slides =========================================

def s01_title(prs):
    s = blank(prs)
    gradient_bg(s)
    text(s, MX, Inches(1.5), SW - 2 * MX, Inches(0.35),
         ["RESEARCH INITIATION PROJECT — FINAL PRESENTATION"],
         size=13, color=LILAC, bold=True)
    text(s, MX, Inches(1.95), Inches(11), Inches(1.6),
         ["Finding Failing Transformers Without Labels"],
         size=44, color=WHITE, bold=True, spacing=1.05)
    text(s, MX, Inches(3.45), Inches(10.5), Inches(0.9),
         ["Label-free health monitoring from dissolved-gas analysis — "
          "and a hidden trap in how such models are validated."],
         size=19, color=LILAC, spacing=1.25)
    text(s, MX, Inches(4.9), Inches(10.5), Inches(1.6), [
        [("Thibaut Faucheux", {"bold": True, "color": WHITE}),
         (" — CESI Engineering School, France", {"color": LILAC})],
        [("Supervised by Dr. Rattanakorn Phadungthin", {"color": WHITE}),
         (" — KMUTNB, Bangkok, Thailand", {"color": LILAC})],
        [("July 2026 · 5-week research project", {"color": LILAC})],
    ], size=15, spacing=1.5)


def s02_agenda(prs):
    s = blank(prs)
    kicker(s, "Plan")
    h2(s, "What the next 20 minutes look like")
    items = [
        ("01", "Context — transformers, gas analysis, and a fleet with no labels"),
        ("02", "Contribution 1 — a fault-type map learned without labels"),
        ("03", "Contribution 2 — the validation trap: surveillance bias"),
        ("04", "Contribution 3 — a fleet risk index + live dashboard demo"),
        ("05", "Take-aways — for the field, and for me"),
    ]
    y = Inches(1.85)
    for num, label in items:
        text(s, MX, y, Inches(0.8), Inches(0.5), [num], size=17, color=ACCENT, bold=True)
        text(s, MX + Inches(0.9), y, Inches(10.8), Inches(0.5), [label], size=19, bold=True)
        ln = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, MX, y + Inches(0.62),
                                    SW - MX, y + Inches(0.62))
        ln.line.color.rgb = BORDER; ln.line.width = Pt(0.75); _noshadow(ln)
        y += Inches(0.92)
    footer(s, 2)


def s03_dga(prs):
    s = blank(prs)
    kicker(s, "1 · Context")
    h2(s, "A blood test for transformers")
    text(s, MX, Inches(1.38), SW - 2 * MX, Inches(0.85),
         ["Power transformers are critical, million-euro assets. When something goes wrong "
          "inside — overheating, sparking, arcing — the insulating oil breaks down and releases "
          "dissolved gases. Measuring them is called DGA: Dissolved Gas Analysis."],
         size=15.5, color=MUTED, spacing=1.3)
    # diagram: three boxes + arrows
    by = Inches(2.7); bh = Inches(2.0)
    rect(s, Inches(0.9), by, Inches(2.6), bh, fill=ACCENT_SOFT, line=ACCENT)
    text(s, Inches(0.9), by + Inches(0.6), Inches(2.6), Inches(0.9),
         [[("Transformer", {"bold": True})], [("fault heats the oil", {"color": MUTED, "size": 12})]],
         size=15, align=PP_ALIGN.CENTER)
    arrow(s, Inches(3.6), by + bh / 2, Inches(4.25), by + bh / 2)
    rect(s, Inches(4.35), by, Inches(2.6), bh)
    for cx, cy, r in [(5.1, 4.15, 0.10), (5.6, 3.75, 0.07), (6.1, 3.35, 0.11),
                      (5.4, 3.15, 0.06), (6.0, 2.95, 0.07)]:
        c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx), Inches(cy),
                               Inches(2 * r), Inches(2 * r))
        c.fill.solid(); c.fill.fore_color.rgb = ACCENT; c.line.fill.background()
        _noshadow(c)
    text(s, Inches(4.35), by + bh - Inches(0.5), Inches(2.6), Inches(0.4),
         [[("Gases dissolve in oil", {"bold": True})]], size=13, align=PP_ALIGN.CENTER)
    arrow(s, Inches(7.05), by + bh / 2, Inches(7.7), by + bh / 2)
    rect(s, Inches(7.8), by, Inches(2.2), bh)
    for bx, bhh, col in [(8.15, 0.55, ACCENT), (8.6, 0.95, CYAN), (9.05, 0.35, FAINT)]:
        b = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(bx),
                               by + bh - Inches(0.55) - Inches(bhh), Inches(0.3), Inches(bhh))
        b.fill.solid(); b.fill.fore_color.rgb = col; b.line.fill.background(); _noshadow(b)
    text(s, Inches(7.8), by + bh - Inches(0.5), Inches(2.2), Inches(0.4),
         [[("Lab: 7 gases (ppm)", {"bold": True})]], size=13, align=PP_ALIGN.CENTER)
    # clue card
    card(s, Inches(10.35), Inches(2.5), Inches(2.5), Inches(2.6), fill=WHITE)
    text(s, Inches(10.55), Inches(2.62), Inches(2.15), Inches(0.35),
         [[("Each gas is a clue:", {"bold": True})]], size=13.5)
    text(s, Inches(10.55), Inches(3.0), Inches(2.15), Inches(2.0), [
        [("H2", {"bold": True}), (" — sparking", {})],
        [("C2H2", {"bold": True}), (" — arcing (worst)", {})],
        [("C2H4/CH4", {"bold": True}), (" — overheating", {})],
        [("CO/CO2", {"bold": True}), (" — paper ageing", {})],
    ], size=12, spacing=1.35)
    footer(s, 3)


def s04_problem(prs):
    s = blank(prs)
    kicker(s, "1 · Context")
    h2(s, "The problem: experts are scarce, labels don't exist")
    pic_fit(s, "duval_triangle", MX, Inches(1.6), Inches(5.4), Inches(4.6))
    text(s, MX, Inches(6.3), Inches(5.4), Inches(0.6),
         ["The Duval Triangle — an expert rule from the 1970s, still the industry standard."],
         size=11, color=FAINT, align=PP_ALIGN.CENTER)
    bx = Inches(6.6); bw = SW - bx - MX
    bullet(s, bx, Inches(1.85), bw, [
        [("Expert rules", {"bold": True}),
         (" (Duval, IEC ratios) work — but need an expert, and often answer "
          "“undetermined”.", {})],
        [("Machine learning", {"bold": True}),
         (" reaches 99% accuracy in papers — but needs labelled faults, which utilities "
          "almost never have.", {})],
        [("Our fleet: zero confirmed fault labels.", {"bold": True}),
         (" Just gas numbers and messy operator notes.", {})],
    ], size=15, gap=1.0)
    card(s, bx, Inches(5.15), bw, Inches(1.35), fill=ACCENT_SOFT, line=ACCENT)
    text(s, bx + Inches(0.25), Inches(5.38), bw - Inches(0.5), Inches(0.95),
         [[("Goal: ", {"bold": True}),
           ("rank 628 transformers by risk and read fault types — ", {}),
           ("using no labels at all.", {"bold": True})]], size=16, spacing=1.25)
    footer(s, 4)


def s05_data(prs):
    s = blank(prs)
    kicker(s, "1 · Context")
    h2(s, "Real fleet data — not a curated benchmark")
    tw = (SW - 2 * MX - 3 * Inches(0.25)) / 4
    vals = [("628", "transformers in service"), ("4,563", "oil samples analysed"),
            ("5 yrs", "of history (2019–2024)"), ("42", "manufacturers")]
    for i, (n, l) in enumerate(vals):
        tile(s, MX + i * (tw + Inches(0.25)), Inches(1.65), tw, Inches(1.25), n, l)
    cw = (SW - 2 * MX - Inches(0.3)) / 2
    card(s, MX, Inches(3.25), cw, Inches(3.3))
    text(s, MX + Inches(0.25), Inches(3.45), cw - Inches(0.5), Inches(0.4),
         [[("What “real” means here", {"bold": True})]], size=16)
    bullet(s, MX + Inches(0.25), Inches(3.95), cw - Inches(0.5), [
        "Gas values stored as text, with dashes for missing",
        "Operator notes mix English and Thai — 69% contain Thai",
        "Notes log everything: faults, but mostly maintenance and admin",
        "No ground truth; assembled from multiple sources, anonymised",
    ], size=13.5, gap=0.62)
    card(s, MX + cw + Inches(0.3), Inches(3.25), cw, Inches(3.3))
    text(s, MX + cw + Inches(0.55), Inches(3.45), cw - Inches(0.6), Inches(0.4),
         [[("Why that is interesting", {"bold": True})]], size=16)
    text(s, MX + cw + Inches(0.55), Inches(3.95), cw - Inches(0.6), Inches(2.4),
         ["Most published methods are tested on small, clean, balanced datasets (often under "
          "200 samples). Whether they survive contact with a real fleet is exactly the open "
          "question — and the messiness itself hides a scientific finding."],
         size=13.5, color=MUTED, spacing=1.35)
    footer(s, 5)


def s06_journey(prs):
    s = blank(prs)
    kicker(s, "1 · Context")
    h2(s, "Five weeks that did not go as planned — on purpose")
    ly = Inches(2.6)
    ln = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.0), ly, SW - Inches(1.0), ly)
    ln.line.color.rgb = BORDER; ln.line.width = Pt(3); _noshadow(ln)
    weeks = [
        (1.45, ACCENT, "Week 1", "Clean the data,\nrebuild expert rules", 0.14),
        (4.05, ACCENT, "Week 2", "First models +\na fleet risk index", 0.14),
        (6.65, BAD, "Week 3 — discovery", "Our validation label is biased.\nEverything re-examined.", 0.18),
        (9.25, ACCENT, "Week 4", "Attack our own results: simpler\nwins; Thai notes translated", 0.14),
        (11.85, GOOD, "Week 5", "IEEE paper, dashboard,\nthis defense", 0.14),
    ]
    for cx, col, wk, desc, r in weeks:
        c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - r), ly - Inches(r),
                               Inches(2 * r), Inches(2 * r))
        c.fill.solid(); c.fill.fore_color.rgb = col; c.line.fill.background(); _noshadow(c)
        text(s, Inches(cx - 1.25), ly - Inches(0.75), Inches(2.5), Inches(0.4),
             [wk], size=13.5, bold=True, align=PP_ALIGN.CENTER,
             color=col if col != ACCENT else INK)
        for j, dline in enumerate(desc.split("\n")):
            text(s, Inches(cx - 1.3), ly + Inches(0.3 + 0.27 * j), Inches(2.6), Inches(0.3),
                 [dline], size=11.5, color=MUTED, align=PP_ALIGN.CENTER)
    card(s, Inches(2.4), Inches(5.2), SW - Inches(4.8), Inches(1.15))
    text(s, Inches(2.7), Inches(5.42), SW - Inches(5.4), Inches(0.8),
         [[("Research is not a straight line. Our two strongest findings came from ", {}),
           ("questioning our own results", {"bold": True}),
           (" — not from the original plan.", {})]],
         size=15, align=PP_ALIGN.CENTER, spacing=1.25)
    footer(s, 6)


def s07_severity_trap(prs):
    s = blank(prs)
    kicker(s, "2 · Contribution 1 — fault types without labels")
    h2(s, "The severity trap")
    text(s, MX, Inches(1.4), SW - 2 * MX, Inches(0.8),
         [[("Feed raw gas concentrations to any model, and it learns ", {}),
           ("“how much gas”", {"bold": True}),
           (" — not ", {}), ("“which fault”", {"bold": True}),
           (". The same fault at low and high gassing looks like two different things.", {})]],
         size=15.5, color=MUTED, spacing=1.3)
    py = Inches(2.55); ph = Inches(2.9)
    pw = Inches(5.4)
    card(s, MX, py, pw, ph, fill=OFFW)
    text(s, MX, py + Inches(0.12), pw, Inches(0.35),
         [[("Concentrations (ppm)", {"bold": True})]], size=14, align=PP_ALIGN.CENTER)
    for cx, cy in [(1.9, 4.6), (5.0, 3.3)]:
        c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx), Inches(cy), Inches(0.32), Inches(0.32))
        c.fill.solid(); c.fill.fore_color.rgb = BAD; c.line.fill.background(); _noshadow(c)
    text(s, Inches(1.35), Inches(4.25), Inches(1.6), Inches(0.3), ["arcing, early"],
         size=11, color=MUTED, align=PP_ALIGN.CENTER)
    text(s, Inches(4.45), Inches(2.95), Inches(1.6), Inches(0.3), ["arcing, severe"],
         size=11, color=MUTED, align=PP_ALIGN.CENTER)
    text(s, MX, py + ph - Inches(0.55), pw, Inches(0.4),
         [[("same fault — far apart", {"bold": True, "color": BAD})]],
         size=13, align=PP_ALIGN.CENTER)
    arrow(s, Inches(6.35), py + ph / 2, Inches(7.0), py + ph / 2, color=FAINT)
    x2 = Inches(7.15)
    card(s, x2, py, pw, ph, fill=OFFW, line=ACCENT)
    text(s, x2, py + Inches(0.12), pw, Inches(0.35),
         [[("Proportions (the gas “recipe”)", {"bold": True})]],
         size=14, align=PP_ALIGN.CENTER)
    for cx, cy, op in [(9.6, 3.75, GOOD), (9.85, 3.95, GOOD)]:
        c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx), Inches(cy), Inches(0.32), Inches(0.32))
        c.fill.solid(); c.fill.fore_color.rgb = op; c.line.fill.background(); _noshadow(c)
    text(s, x2, py + ph - Inches(0.55), pw, Inches(0.4),
         [[("same fault — same place", {"bold": True, "color": GOOD})]],
         size=13, align=PP_ALIGN.CENTER)
    expert(s, MX, Inches(5.85), SW - 2 * MX, Inches(0.95),
           "A plain autoencoder on log-concentrations encodes magnitude: total gas is linearly "
           "readable from its 2-D latent (R² = 0.63), and agreement with Duval classes is "
           "near-noise (ARI 0.14–0.16).")
    footer(s, 7)


def s08_c1_result(prs):
    s = blank(prs)
    kicker(s, "2 · Contribution 1 — fault types without labels")
    h2(s, "Use the right geometry, and fault types appear")
    pic_fit(s, "representation_ablation", MX, Inches(1.65), Inches(7.0), Inches(4.5))
    text(s, MX, Inches(6.25), Inches(7.0), Inches(0.5),
         ["Agreement with the expert (Duval) classes, from raw inputs to our representation."],
         size=11, color=FAINT, align=PP_ALIGN.CENTER)
    bx = Inches(8.0); bw = SW - bx - MX
    rect(s, bx, Inches(1.65), bw, Inches(1.3))
    text(s, bx + Inches(0.2), Inches(1.78), bw - Inches(0.4), Inches(0.55),
         [[("0.16 → 0.55", {"bold": True, "color": ACCENT, "size": 28})]])
    text(s, bx + Inches(0.2), Inches(2.38), bw - Inches(0.4), Inches(0.5),
         [[("agreement with expert diagnosis (ARI) — ", {}),
           ("3.4× better", {"bold": True}), (", with zero labels", {})]],
         size=12, color=MUTED)
    card(s, bx, Inches(3.15), bw, Inches(1.35))
    text(s, bx + Inches(0.2), Inches(3.32), bw - Inches(0.4), Inches(1.05),
         [[("The trick is not a bigger model. ", {"bold": True}),
           ("It is treating gas proportions with the mathematics made for proportions "
            "(log-ratios), then the simplest possible projection.", {})]],
         size=13.5, spacing=1.25)
    expert(s, bx, Inches(4.75), bw, Inches(1.75),
           "Centred log-ratio (Aitchison) on the 5 combustible gases → standardise → "
           "PCA-2 → KMeans (k=7). ARI 0.545 ± 0.002 vs Duval. Deterministic. Holds on "
           "a temporal split: fit pre-2022, score 2022+ → ARI 0.45.")
    footer(s, 8)


def s09_map(prs):
    s = blank(prs)
    kicker(s, "2 · Contribution 1 — fault types without labels")
    h2(s, "A data-driven diagnostic map — and an honest negative")
    pic_fit(s, "clr_pca_map", MX, Inches(1.6), Inches(6.4), Inches(4.7))
    text(s, MX, Inches(6.4), Inches(6.4), Inches(0.55),
         ["Every dot is an oil sample. Colours = expert diagnosis, added afterwards — "
          "the map itself never saw them."], size=11, color=FAINT, align=PP_ALIGN.CENTER)
    bx = Inches(7.45); bw = SW - bx - MX
    bullet(s, bx, Inches(1.85), bw, [
        [("Discharge, partial-discharge and thermal regions ", {}),
         ("emerge on their own", {"bold": True})],
        "Readable axes — an engineer can interpret the directions",
        "Deterministic: same data in, same map out",
    ], size=14, gap=0.75)
    card(s, bx, Inches(4.35), bw, Inches(2.15), fill=WHITE, line=BAD)
    text(s, bx + Inches(0.22), Inches(4.52), bw - Inches(0.44), Inches(0.35),
         [[("NEGATIVE RESULT", {"bold": True, "color": BAD, "size": 11})]])
    text(s, bx + Inches(0.22), Inches(4.9), bw - Inches(0.44), Inches(1.5),
         [[("We also built a deep autoencoder (SD-CAE) and an adversarial variant. ", {}),
           ("Neither beat the simple linear map. ", {"bold": True}),
           ("We report both as negatives — knowing what does not work is a result too.", {})]],
         size=13, spacing=1.3)
    footer(s, 9)


def s10_loop(prs):
    s = blank(prs)
    kicker(s, "3 · Contribution 2 — the validation trap")
    h2(s, "“The more you look, the more you find”")
    text(s, MX, Inches(1.4), SW - 2 * MX, Inches(0.6),
         ["To check our risk score, we used the only “truth” available: "
          "operator-logged field events. Then we noticed something wrong."],
         size=15, color=MUTED, spacing=1.3)
    by = Inches(2.7); bh = Inches(1.3); bw = Inches(2.5); gap = Inches(0.55)
    boxes = [
        ("Operator worries", "about a transformer", ACCENT_SOFT, ACCENT, INK),
        ("Samples it more", "extra oil tests", WHITE, BORDER, INK),
        ("Logs more notes", "events get recorded", WHITE, BORDER, INK),
        ("“Risky” in the data", "the loop feeds itself", RGBColor(0xFD, 0xEA, 0xEA), BAD, BAD),
    ]
    x = Inches(0.9)
    centers = []
    for t1, t2, fill, line, tc in boxes:
        rect(s, x, by, bw, bh, fill=fill, line=line)
        text(s, x, by + Inches(0.22), bw, Inches(0.45), [[(t1, {"bold": True, "color": tc})]],
             size=15, align=PP_ALIGN.CENTER)
        text(s, x, by + Inches(0.68), bw, Inches(0.4), [t2], size=12, color=MUTED,
             align=PP_ALIGN.CENTER)
        centers.append(x + bw / 2)
        if x + bw + gap < SW - Inches(0.9):
            arrow(s, x + bw + Inches(0.05), by + bh / 2, x + bw + gap - Inches(0.05), by + bh / 2)
        x += bw + gap
    # feedback arrow (elbow)
    fb = s.shapes.add_connector(MSO_CONNECTOR.ELBOW, centers[-1], by + bh,
                                centers[0], by + bh + Inches(0.9))
    _arrowhead(fb, BAD)
    fb.line.dash_style = None
    _noshadow(fb)
    text(s, Inches(2.5), by + bh + Inches(1.0), SW - Inches(5.0), Inches(0.4),
         [[("the loop feeds itself — attention, not chemistry", {"bold": True, "color": BAD})]],
         size=14, align=PP_ALIGN.CENTER)
    expert(s, MX, Inches(5.85), SW - 2 * MX, Inches(1.0),
           "This is surveillance / informed-presence bias, well known in medicine and health "
           "records (Sackett 1979; Haut & Pronovost 2011; Goldstein 2016). To our knowledge, "
           "never reported for DGA — models in this field are routinely validated against such labels.")
    footer(s, 10)


def s11_evidence(prs):
    s = blank(prs)
    kicker(s, "3 · Contribution 2 — the validation trap")
    h2(s, "Three tests, one verdict")
    pic_fit(s, "confound_panel", MX, Inches(1.55), SW - 2 * MX, Inches(3.4))
    cw = (SW - 2 * MX - Inches(0.5)) / 3
    texts = [
        [("(a) The tie. ", {"bold": True}),
         ("Counting how often a unit was sampled predicts logged events as well as all our "
          "chemistry combined (0.76 vs 0.74 — statistically indistinguishable).", {})],
        [("(b) No foresight. ", {"bold": True}),
         ("Control for the sample count, and the score's ability to predict future events "
          "drops to a coin flip (0.50).", {})],
        [("(c) The fix. ", {"bold": True}),
         ("Define the target by chemistry (arcing onset) instead of attention — now gases "
          "predict (0.64) and the count does not (0.49).", {})],
    ]
    for i, t in enumerate(texts):
        xx = MX + i * (cw + Inches(0.25))
        card(s, xx, Inches(5.15), cw, Inches(1.75))
        text(s, xx + Inches(0.18), Inches(5.32), cw - Inches(0.36), Inches(1.45),
             [t], size=12.5, spacing=1.25)
    footer(s, 11)


def s12_thai(prs):
    s = blank(prs)
    kicker(s, "3 · Contribution 2 — the validation trap")
    h2(s, "Repairing the label: 211 Thai notes, read one by one")
    tw = (SW - 2 * MX - 2 * Inches(0.25)) / 3
    tiles = [("69%", "of operator notes contain Thai — invisible to the English keyword filter"),
             ("211", "distinct Thai notes translated and classified (released for expert validation)"),
             ("+6", "genuine electrical events recovered (noise, breaker explosions, insulation symptoms)")]
    for i, (n, l) in enumerate(tiles):
        tile(s, MX + i * (tw + Inches(0.25)), Inches(1.65), tw, Inches(1.5), n, l,
             num_color=ACCENT)
    cw = (SW - 2 * MX - Inches(0.3)) / 2
    card(s, MX, Inches(3.5), cw, Inches(2.9))
    text(s, MX + Inches(0.25), Inches(3.68), cw - Inches(0.5), Inches(0.4),
         [[("Example finds", {"bold": True})]], size=15.5)
    text(s, MX + Inches(0.25), Inches(4.15), cw - Inches(0.5), Inches(2.1), [
        [("“หม้อแปลงมีเสียงดัง”", {}),
         (" — transformer making abnormal noise", {"italic": True, "color": MUTED})],
        [("“ตามผลเนื่องจาก Bkr. ระเบิด”", {}),
         (" — follow-up after breaker explosion", {"italic": True, "color": MUTED})],
        [("“ค่า p.f สูง”", {}),
         (" — high insulation power factor", {"italic": True, "color": MUTED})],
    ], size=13.5, spacing=1.6)
    card(s, MX + cw + Inches(0.3), Inches(3.5), cw, Inches(2.9), fill=ACCENT_SOFT, line=ACCENT)
    text(s, MX + cw + Inches(0.55), Inches(3.68), cw - Inches(0.6), Inches(0.4),
         [[("The punchline", {"bold": True})]], size=15.5)
    text(s, MX + cw + Inches(0.55), Inches(4.15), cw - Inches(0.6), Inches(2.1),
         [[("On the repaired label, the sample-count's advantage ", {}),
           ("disappears entirely", {"bold": True}),
           (" (0.734 vs 0.735). Part of the “trap” was simply that the label was "
            "incomplete — in a language the pipeline didn't read.", {})]],
         size=14, spacing=1.4)
    footer(s, 12)


def s13_index(prs):
    s = blank(prs)
    kicker(s, "4 · Contribution 3 — from research to practice")
    h2(s, "A triage list for the maintenance team")
    pic_fit(s, "health_decile_note_rate", MX, Inches(1.6), Inches(7.2), Inches(4.5))
    text(s, MX, Inches(6.2), Inches(7.2), Inches(0.6),
         ["Fleet sorted into 10 risk groups by our label-free index. Real logged events "
          "concentrate exactly where the index says they should."],
         size=11, color=FAINT, align=PP_ALIGN.CENTER)
    bx = Inches(8.2); bw = SW - bx - MX
    tile(s, bx, Inches(1.6), bw, Inches(1.15), "27%",
         "event rate in the riskiest 10% of the fleet…", num_color=BAD)
    tile(s, bx, Inches(2.95), bw, Inches(1.15), "2%",
         "…versus the safest 10%", num_color=GOOD)
    card(s, bx, Inches(4.35), bw, Inches(2.1))
    text(s, bx + Inches(0.2), Inches(4.52), bw - Inches(0.4), Inches(1.8),
         [[("Four transparent ingredients: ", {"bold": True}),
           ("gas severity, acetylene (arcing), the speed of hydrogen rise, and statistical "
            "novelty. No black box — and always reported together with the surveillance-bias "
            "floor from the previous slides.", {})]],
         size=13, spacing=1.3)
    footer(s, 13)


def s14_demo(prs):
    s = blank(prs)
    kicker(s, "4 · Contribution 3 — from research to practice")
    text(s, MX, Inches(1.7), SW - 2 * MX, Inches(1.2), ["DEMO"],
         size=72, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    text(s, MX, Inches(2.9), SW - 2 * MX, Inches(0.6), ["The whole project, live"],
         size=26, bold=True, align=PP_ALIGN.CENTER)
    cw = Inches(8.4); cx = (SW - cw) / 2
    card(s, cx, Inches(3.7), cw, Inches(2.15))
    text(s, cx + Inches(0.45), Inches(3.92), cw - Inches(0.9), Inches(1.8), [
        [("1.  ", {"bold": True, "color": ACCENT}),
         ("Fleet ranking — 628 units re-ranked live as we switch scoring profile", {})],
        [("2.  ", {"bold": True, "color": ACCENT}),
         ("Transformer inspector — gas history, expert diagnoses, “why this score”", {})],
        [("3.  ", {"bold": True, "color": ACCENT}),
         ("The fault-type map and the surveillance-bias evidence — the paper, interactive", {})],
    ], size=14.5, spacing=1.6)
    text(s, MX, Inches(6.1), SW - 2 * MX, Inches(0.5),
         ["localhost:8501 — runs on the anonymised dataset · if the demo gods object, "
          "every view also exists as a paper figure"],
         size=12, color=MUTED, align=PP_ALIGN.CENTER)
    footer(s, 14)


def s15_takeaways(prs):
    s = blank(prs)
    kicker(s, "5 · Take-aways")
    h2(s, "What this project says to the field")
    cw = (SW - 2 * MX - Inches(0.5)) / 3
    cards = [
        ("1", "The right geometry beats a bigger model.",
         "A linear method in the correct mathematical space outperformed deep networks — "
         "simpler, interpretable, deployable."),
        ("2", "Validation labels can lie.",
         "Any DGA risk model validated on operator-logged events should report the "
         "sample-count “confound floor” — otherwise surveillance can be mistaken "
         "for diagnostic skill."),
        ("3", "Honest negatives are results.",
         "Deep SVDD, VAE, adversarial disentanglement, RUL modelling: tested, rejected, "
         "documented. Future work starts further ahead."),
    ]
    for i, (n, t, b) in enumerate(cards):
        xx = MX + i * (cw + Inches(0.25))
        card(s, xx, Inches(1.7), cw, Inches(3.15))
        badge = s.shapes.add_shape(MSO_SHAPE.OVAL, xx + Inches(0.22), Inches(1.92),
                                   Inches(0.42), Inches(0.42))
        badge.fill.solid(); badge.fill.fore_color.rgb = ACCENT; badge.line.fill.background()
        _noshadow(badge)
        text(s, xx + Inches(0.22), Inches(1.955), Inches(0.42), Inches(0.36), [n],
             size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        text(s, xx + Inches(0.22), Inches(2.55), cw - Inches(0.44), Inches(0.7),
             [[(t, {"bold": True})]], size=15, spacing=1.15)
        text(s, xx + Inches(0.22), Inches(3.3), cw - Inches(0.44), Inches(1.45),
             [b], size=12, color=MUTED, spacing=1.3)
    card(s, MX, Inches(5.3), SW - 2 * MX, Inches(1.25))
    text(s, MX + Inches(0.3), Inches(5.5), SW - 2 * MX - Inches(0.6), Inches(0.9),
         [[("Deliverables:  ", {"bold": True}),
           ("a 6-page IEEE paper draft (3 contributions) · a reproducible pipeline (one command "
            "regenerates every number, 10 automated tests) · an anonymised public dataset "
            "extract · the interactive dashboard.", {})]],
         size=13.5, spacing=1.3)
    footer(s, 15)


def s16_personal(prs):
    s = blank(prs)
    kicker(s, "5 · Take-aways")
    h2(s, "What five weeks of research taught me")
    cw = (SW - 2 * MX - Inches(0.5)) / 2
    left = [
        [("The scientific method, for real", {"bold": True}),
         (" — hypothesis, test, and the discipline to pivot when the evidence disagrees "
          "with the plan", {})],
        [("Statistical honesty", {"bold": True}),
         (" — confidence intervals, paired tests, blocked bootstraps; “0.76 vs 0.74” "
          "means nothing until you test it", {})],
        [("Attack your own results first", {"bold": True}),
         (" — the finding that a simple method beat our own headline model came from auditing "
          "ourselves before reviewers would", {})],
    ]
    right = [
        [("A full research stack", {"bold": True}),
         (" — Python/PyTorch, statistics, LaTeX/IEEE writing, git, automated testing, "
          "Streamlit", {})],
        [("Working across languages and cultures", {"bold": True}),
         (" — a French student, a Thai fleet, English science; 211 Thai field notes read "
          "and classified", {})],
        [("Scoping", {"bold": True}),
         (" — five weeks forces choices: what to cut, what to consolidate, what to hand to "
          "future work", {})],
    ]
    bullet(s, MX, Inches(2.0), cw, left, size=14.5, gap=1.5)
    bullet(s, MX + cw + Inches(0.5), Inches(2.0), cw, right, size=14.5, gap=1.5)
    footer(s, 16)


def s17_thanks(prs):
    s = blank(prs)
    gradient_bg(s)
    text(s, MX, Inches(1.9), SW - 2 * MX, Inches(0.35), ["THANK YOU"],
         size=13, color=LILAC, bold=True)
    text(s, MX, Inches(2.35), SW - 2 * MX, Inches(1.1), ["Questions?"],
         size=48, color=WHITE, bold=True)
    text(s, MX, Inches(3.7), Inches(10.5), Inches(1.0), [
        "Thibaut Faucheux · CESI × KMUTNB",
        "With thanks to Dr. Rattanakorn Phadungthin for the supervision and domain guidance.",
    ], size=16, color=LILAC, spacing=1.5)
    text(s, MX, Inches(5.3), SW - 2 * MX, Inches(0.8),
         ["628 transformers · 4,563 samples · ARI 0.16 → 0.55 without labels · "
          "one validation trap exposed · 6-page IEEE paper · live dashboard"],
         size=13, color=LILAC)


def main():
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    for build in [s01_title, s02_agenda, s03_dga, s04_problem, s05_data, s06_journey,
                  s07_severity_trap, s08_c1_result, s09_map, s10_loop, s11_evidence,
                  s12_thai, s13_index, s14_demo, s15_takeaways, s16_personal, s17_thanks]:
        build(prs)
    prs.save(OUT)
    print(f"wrote {OUT} ({OUT.stat().st_size/1024:.0f} KB, {len(prs.slides._sldIdLst)} slides)")


if __name__ == "__main__":
    main()
